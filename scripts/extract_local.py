#!/usr/bin/env python3
"""Extract text from local PDFs, DOCX, and PPTX files referenced in _resources.

Walks site/_resources/ for Markdown files with a `local_path` front-matter field,
extracts readable text from the referenced file, and writes it to
site/_search/cache/local/<slug>.txt.

Idempotent: re-runs overwrite only if the source file is newer than the cached text.

Usage:
    python scripts/extract_local.py [--force]
"""

import argparse
import os
import re
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Lazy imports for extraction libraries – fail gracefully with clear message
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path) -> str:
    from pdfminer.high_level import extract_text
    return extract_text(str(path))


def _extract_docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
}

# ---------------------------------------------------------------------------
# YAML front-matter parser (minimal, avoids heavy deps)
# ---------------------------------------------------------------------------

_FM_RE = re.compile(r"^---\s*\n(.*?\n)---\s*\n", re.DOTALL)


def _parse_front_matter(text: str) -> dict:
    """Return a dict of simple scalar front-matter values."""
    m = _FM_RE.match(text)
    if not m:
        return {}
    result = {}
    for line in m.group(1).splitlines():
        if ":" in line and not line.startswith(" ") and not line.startswith("-"):
            key, _, val = line.partition(":")
            val = val.strip().strip('"').strip("'")
            if val:
                result[key.strip()] = val
    return result


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--force", action="store_true",
                        help="Re-extract even if cache is newer than source")
    args = parser.parse_args()

    repo_root = Path(__file__).resolve().parent.parent
    resources_dir = repo_root / "site" / "_resources"
    cache_dir = repo_root / "site" / "_search" / "cache" / "local"
    cache_dir.mkdir(parents=True, exist_ok=True)

    if not resources_dir.is_dir():
        print(f"ERROR: resources directory not found: {resources_dir}", file=sys.stderr)
        return 1

    extracted = 0
    skipped = 0
    errors = 0

    for md_file in sorted(resources_dir.rglob("*.md")):
        fm = _parse_front_matter(md_file.read_text(encoding="utf-8"))
        local_path_str = fm.get("local_path")
        if not local_path_str:
            continue

        source_path = repo_root / local_path_str
        if not source_path.is_file():
            print(f"WARNING: local_path not found: {source_path} (from {md_file.name})",
                  file=sys.stderr)
            errors += 1
            continue

        slug = md_file.stem
        cache_file = cache_dir / f"{slug}.txt"

        # Skip if cache is newer than source (unless --force)
        if (not args.force
                and cache_file.exists()
                and cache_file.stat().st_mtime >= source_path.stat().st_mtime):
            skipped += 1
            continue

        ext = source_path.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if extractor is None:
            print(f"WARNING: unsupported file type '{ext}' for {source_path}",
                  file=sys.stderr)
            errors += 1
            continue

        try:
            text = extractor(source_path)
            cache_file.write_text(text, encoding="utf-8")
            extracted += 1
            print(f"  extracted: {slug} ({len(text)} chars)")
        except Exception as exc:
            print(f"ERROR: failed to extract {source_path}: {exc}", file=sys.stderr)
            errors += 1

    print(f"\nDone: {extracted} extracted, {skipped} skipped (cached), {errors} errors")
    return 1 if errors > 0 and extracted == 0 else 0


if __name__ == "__main__":
    sys.exit(main())
